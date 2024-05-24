"""
作者：吕枫烨
创建时间：2024/4/11 17:25
项目：Python-Core-50-Courses
包：
标题：
描述：
"""

print('hello, world')

# region    # ------------ 判断变量类型 ------------
a = 100
b = 12.345
c = 'hello, world'
d = True
print(type(a))    # <class 'int'>
print(type(b))    # <class 'float'>
print(type(c))    # <class 'str'>
print(type(d))    # <class 'bool'>
# endregion

# region    # ------------ 变量类型转换 ------------
a = 100
b = 12.345
c = 'hello, world'
d = True
# 整数转成浮点数
print(float(a))    # 100.0
# 浮点型转成字符串 (输出字符串时不会看到引号哟)
print(str(b))      # 12.345
# 字符串转成布尔型 (有内容的字符串都会变成True)
print(bool(c))     # True
# 布尔型转成整数 (True会转成1，False会转成0)
print(int(d))      # 1
# 将整数变成对应的字符 (97刚好对应字符表中的字母a)
print(chr(97))     # a
# 将字符转成整数 (Python中字符和字符串表示法相同)
print(ord('a'))    # 97
# endregion

# region    # ------------ 算术运算 ------------

print(321 + 123)     # 加法运算
print(321 - 123)     # 减法运算
print(321 * 123)     # 乘法运算
print(321 / 123)     # 除法运算
print(321 % 123)     # 求模运算
print(321 // 123)    # 整除运算
print(321 ** 123)    # 求幂运算
# endregion

# region    # ------------ 赋值运算符 ------------

a = 10
b = 3
a += b        # 相当于：a = a + b
a *= a + 2    # 相当于：a = a * (a + 2)
print(a)      # 算一下这里会输出什么
# endregion

# region    # ------------ 比较运算符和逻辑运算符 ------------

flag0 = 1 == 1
flag1 = 3 > 2
flag2 = 2 < 1
flag3 = flag1 and flag2
flag4 = flag1 or flag2
flag5 = not (1 != 2)
print('flag0 =', flag0)    # flag0 = True
print('flag1 =', flag1)    # flag1 = True
print('flag2 =', flag2)    # flag2 = False
print('flag3 =', flag3)    # flag3 = False
print('flag4 =', flag4)    # flag4 = True
print('flag5 =', flag5)    # flag5 = False
# endregion

# region    例子1：华氏温度转换为摄氏温度。
f = float(input('请输入华氏温度: '))
c = (f - 32) / 1.8
print('%.1f华氏度 = %.1f摄氏度' % (f, c))
print(f'{f:.1f}华氏度 = {c:.1f}摄氏度')
# endregion

# region    例子2：输入圆的半径计算计算周长和面积。
radius = float(input('请输入圆的半径: '))
perimeter = 2 * 3.1416 * radius
area = 3.1416 * radius * radius
print('周长: %.2f' % perimeter)
print('面积: %.2f' % area)
# endregion

# region    例子3：输入年份判断是不是闰年。
year = int(input('请输入年份: '))
is_leap = year % 4 == 0 and year % 100 != 0 or year % 400 == 0
print(is_leap)
# endregion

# region    # ------------ if语句的使用 ------------

username = input('请输入用户名: ')
password = input('请输入口令: ')
# 用户名是admin且密码是123456则身份验证成功否则身份验证失败
if username == 'admin' and password == '123456':
    print('身份验证成功!')
else:
    print('身份验证失败!')

x = float(input('x = '))
if x > 1:
    y = 3 * x - 5
elif x >= -1:
    y = x + 2
else:
    y = 5 * x + 3
print(f'f({x}) = {y}')
# endregion

# region    # 例子1：英制单位英寸与公制单位厘米互换。

value = float(input('请输入长度: '))
unit = input('请输入单位: ')
if unit == 'in' or unit == '英寸':
    print('%f英寸 = %f厘米' % (value, value * 2.54))
elif unit == 'cm' or unit == '厘米':
    print('%f厘米 = %f英寸' % (value, value / 2.54))
else:
    print('请输入有效的单位')
# endregion

# region    # 例子2：百分制成绩转换为等级制成绩。
# score = float(input('请输入成绩: '))
score = 90
if score >= 90:
    grade = 'A'
elif score >= 80:
    grade = 'B'
elif score >= 70:
    grade = 'C'
elif score >= 60:
    grade = 'D'
else:
    grade = 'E'
print('对应的等级是:', grade)
# endregion

# region    # 例子3：输入三条边长，如果能构成三角形就计算周长和面积。

# a = float(input('a = '))
# b = float(input('b = '))
# c = float(input('c = '))
a = float(1.2)
b = float(1.5)
c = float(1)
if a + b > c and a + c > b and b + c > a:
    peri = a + b + c
    print(f'周长: {peri:.2f}')
    half = peri / 2
    area = (half * (half - a) * (half - b) * (half - c)) ** 0.5
    print(f'面积: {area:.2f}')
else:
    print('不能构成三角形')
# endregion

# region    # ------------ for-in循环 ------------

total = 0
for x in range(1, 101):
    total += x
print(total)

total = 0
for x in range(2, 101, 2):
    total += x
print(total)
# endregion

# region    # ------------ while循环 ------------

import random

# 产生一个1-100范围的随机数
answer = random.randint(1, 100)
counter = 0
while True:
    counter += 1
    number = int(input('请输入: '))
    if number < answer:
        print('大一点')
    elif number > answer:
        print('小一点')
    else:
        print('恭喜你猜对了!')
        break
# 当退出while循环的时候显示用户一共猜了多少次
print(f'你总共猜了{counter}次')
# endregion

# region    # ------------ 嵌套的循环结构 ------------

for i in range(1, 10):
    for j in range(1, i + 1):
        print(f'{i}*{j}={i * j}', end='\t')
    print()
# endregion

# region    例子1：输入一个正整数判断它是不是素数。

num = int(input('请输入一个正整数: '))
end = int(num ** 0.5)
is_prime = True
for x in range(2, end + 1):
    if num % x == 0:
        is_prime = False
        break
if is_prime and num != 1:
    print(f'{num}是素数')
else:
    print(f'{num}不是素数')
# endregion

# region    例子2：输入两个正整数，计算它们的最大公约数和最小公倍数。

# x = int(input('x = '))
# y = int(input('y = '))
x = int(3)
y = int(9)
for factor in range(x, 0, -1):
    if x % factor == 0 and y % factor == 0:
        print(f'{x}和{y}的最大公约数是{factor}')
        print(f'{x}和{y}的最小公倍数是{x * y // factor}')
        break
# endregion

# region    # ------------ 列表 ------------
items1 = list(range(1, 10))
print(items1)    # [1, 2, 3, 4, 5, 6, 7, 8, 9]
items2 = list('hello')
print(items2)    # ['h', 'e', 'l', 'l', 'o']

items1 = [35, 12, 99, 68, 55, 87]
items2 = [45, 8, 29]
# endregion

# region    # ------------ 列表的运算符 ------------
# 列表的拼接
items3 = items1 + items2
print(items3)    # [35, 12, 99, 68, 55, 87, 45, 8, 29]

# 列表的重复
items4 = ['hello'] * 3
print(items4)    # ['hello', 'hello', 'hello']

# 列表的成员运算
print(100 in items3)        # False
print('hello' in items4)    # True

# 获取列表的长度(元素个数)
size = len(items3)
print(size)                 # 9

# 列表的索引
print(items3[0], items3[-size])        # 35 35
items3[-1] = 100
print(items3[size - 1], items3[-1])    # 100 100

# 列表的切片
print(items3[:5])          # [35, 12, 99, 68, 55]
print(items3[4:])          # [55, 87, 45, 8, 100]
print(items3[-5:-7:-1])    # [55, 68]
print(items3[::-2])        # [100, 45, 55, 99, 35]

# 列表的比较运算
items5 = [1, 2, 3, 4]
items6 = list(range(1, 5))
# 两个列表比较相等性比的是对应索引位置上的元素是否相等
print(items5 == items6)    # True
items7 = [3, 2, 1]
# 两个列表比较大小比的是对应索引位置上的元素的大小
print(items5 <= items7)    # True
#endregion

# region    # ------------ 列表元素的遍历 ------------
### 方法一
items = ['Python', 'Java', 'Go', 'Kotlin']

for index in range(len(items)):
    print(items[index])

### 方法二
items = ['Python', 'Java', 'Go', 'Kotlin']

for item in items:
    print(item)
#endregion

# region    案例：掷色子统计每个点数出现次数
import random

counters = [0] * 6
for _ in range(6000):
    face = random.randint(1, 6)
    counters[face - 1] += 1
for face in range(1, 7):
    print(f'{face}点出现了{counters[face - 1]}次')
#endregion

# region    # ------------ 列表的方法：添加和删除元素 ------------
items = ['Python', 'Java', 'Go', 'Kotlin']

# 使用append方法在列表尾部添加元素
items.append('Swift')
print(items)    # ['Python', 'Java', 'Go', 'Kotlin', 'Swift']
# 使用insert方法在列表指定索引位置插入元素
items.insert(2, 'SQL')
print(items)    # ['Python', 'Java', 'SQL', 'Go', 'Kotlin', 'Swift']

# 删除指定的元素
items.remove('Java')
print(items)    # ['Python', 'SQL', 'Go', 'Kotlin', 'Swift']
# 删除指定索引位置的元素，并返回删除的元素
items.pop(0)
items.pop(len(items) - 1)
print(items)    # ['SQL', 'Go', 'Kotlin']
# 删除指定索引位置的元素，并返回删除后的列表，del的删除性能略优
items = ['Python', 'Java', 'Go', 'Kotlin']
del items[1]
print(items)    # ['Python', 'Go', 'Kotlin']

# 清空列表中的元素
items.clear()
print(items)    # []
#endregion

# region    # ------------ 列表的方法：元素位置和次数 ------------
items = ['Python', 'Java', 'Java', 'Go', 'Kotlin', 'Python']
# 查找元素的索引位置
print(items.index('Python'))       # 0
print(items.index('Python', 2))    # 5
# 注意：虽然列表中有'Java'，但是从索引为3这个位置开始后面是没有'Java'的
print(items.index('Java', 3))      # ValueError: 'Java' is not in list

items = ['Python', 'Java', 'Java', 'Go', 'Kotlin', 'Python']
# 查找元素出现的次数
print(items.count('Python'))    # 2
print(items.count('Go'))        # 1
print(items.count('Swfit'))     # 0
#endregion

# region    # ------------ 列表的方法：元素排序和反转 ------------
items = ['Python', 'Java', 'Go', 'Kotlin', 'Python']

# 排序
items.sort()
print(items)    # ['Go', 'Java', 'Kotlin', 'Python', 'Python']
# 反转
items.reverse()
print(items)    # ['Python', 'Python', 'Kotlin', 'Java', 'Go']
#endregion

# region    # ------------ 列表的方法：列表的生成式 ------------
### 方法一：for循环
# 创建一个由1到9的数字构成的列表
items1 = []
for x in range(1, 10):
    items1.append(x)
print(items1)

# 创建一个由'hello world'中除空格和元音字母外的字符构成的列表
items2 = []
for x in 'hello world':
    if x not in ' aeiou':
        items2.append(x)
print(items2)

# 创建一个由个两个字符串中字符的笛卡尔积构成的列表
items3 = []
for x in 'ABC':
    for y in '12':
        items3.append(x + y)
print(items3)

### 方法二：生成式（性能更优）
# 创建一个由1到9的数字构成的列表
items1 = [x for x in range(1, 10)]
print(items1)    # [1, 2, 3, 4, 5, 6, 7, 8, 9]

# 创建一个由'hello world'中除空格和元音字母外的字符构成的列表
items2 = [x for x in 'hello world' if x not in ' aeiou']
print(items2)    # ['h', 'l', 'l', 'w', 'r', 'l', 'd']

# 创建一个由个两个字符串中字符的笛卡尔积构成的列表
items3 = [x + y for x in 'ABC' for y in '12']
print(items3)    # ['A1', 'A2', 'B1', 'B2', 'C1', 'C2']
#endregion

# region    # ------------ 列表的方法：嵌套的列表 ------------
### 错误方法
scores = [[0] * 3] * 5
print(scores)    # [[0, 0, 0], [0, 0, 0], [0, 0, 0], [0, 0, 0], [0, 0, 0]]

scores[0][0] = 95
print(scores)

### 正确方法
scores = [[0] * 3 for _ in range(5)]
scores[0][0] = 95
print(scores)
# [[95, 0, 0], [0, 0, 0], [0, 0, 0], [0, 0, 0], [0, 0, 0]]
#endregion

# region    # ------------ 元组：定义和使用 ------------
# 定义一个三元组
t1 = (30, 10, 55)
# 定义一个四元组
t2 = ('骆昊', 40, True, '四川成都')

# 查看变量的类型
print(type(t1), type(t2))    # <class 'tuple'> <class 'tuple'>
# 查看元组中元素的数量
print(len(t1), len(t2))      # 3 4

# 通过索引运算获取元组中的元素
print(t1[0], t1[-3])         # 30 30
print(t2[3], t2[-1])         # 四川成都 四川成都

# 循环遍历元组中的元素
for member in t2:
    print(member)

# 成员运算
print(100 in t1)    # False
print(40 in t2)     # True

# 拼接
t3 = t1 + t2
print(t3)           # (30, 10, 55, '骆昊', 40, True, '四川成都')

# 切片
print(t3[::3])      # (30, '骆昊', '四川成都')

# 比较运算
print(t1 == t3)    # False
print(t1 >= t3)    # False
print(t1 < (30, 11, 55))    # True
#endregion

# region    # ------------ 元组应用场景：打包 ------------
# 打包
a = 1, 10, 100
b = [1, 10, 100]
print(type(a), a)    # <class 'tuple'> (1, 10, 100)
print(type(b), b)    # <class 'tuple'> (1, 10, 100)
# 解包
i, j, k = a
o, p, q = b
print(i, j, k)       # 1 10 100
print(o, p, q)       # 1 10 100

a = [1, 10, 100, 1000]
i, j, *k = a
print(i, j, k)          # 1 10 [100, 1000]
i, *j, k = a
print(i, j, k)          # 1 [10, 100] 1000
*i, j, k = a
print(i, j, k)          # [1, 10] 100 1000
*i, j = a
print(i, j)             # [1, 10, 100] 1000
i, *j = a
print(i, j)             # 1 [10, 100, 1000]
i, j, k, *l = a
print(i, j, k, l)       # 1 10 100 [1000]
i, j, k, l, *m = a
print(i, j, k, l, m)    # 1 10 100 1000 []

a, b, *c = range(1, 10)
print(a, b, c)
a, b, c = [1, 10, 100]
print(a, b, c)
a, *b, c = 'hello'
print(a, b, c)
#endregion

# region    # ------------ 元组应用场景：交换两个变量的值 ------------
a, b = b, a
a, b, c = b, c, a
#endregion

# region    # ------------ 元组和列表的比较 ------------
import sys
import timeit

a = list(range(100000))
b = tuple(range(100000))
print(sys.getsizeof(a), sys.getsizeof(b))    # 900120 800056

print(timeit.timeit('[1, 2, 3, 4, 5, 6, 7, 8, 9]'))
print(timeit.timeit('(1, 2, 3, 4, 5, 6, 7, 8, 9)'))
#endregion

# region    # ------------ 字符串 ------------
s1 = 'hello, world!'
s2 = "你好，世界！"
print(s1, s2)
# 以三个双引号或单引号开头的字符串可以折行
s3 = '''
hello, 
world!
'''
print(s3)
print(s3, end='')

s1 = '\'hello, world!\''
print(s1)
s2 = '\\hello, world!\\'
print(s2)

# 字符串s1中\t是制表符，\n是换行符
s1 = '\time up \now'
print(s1)
# 字符串s2中没有转义字符，每个字符都是原始含义
s2 = r'\time up \now'
print(s2)

s1 = '\141\142\143\x61\x62\x63'
s2 = '\u9a86\u660a'
print(s1, s2)
#endregion

# region    # ------------ 字符串的运算：拼接和重复 ------------
s1 = 'hello' + ' ' + 'world'
print(s1)    # hello world
s2 = '!' * 3
print(s2)    # !!!
s1 += s2     # s1 = s1 + s2
print(s1)    # hello world!!!
s1 *= 2      # s1 = s1 * 2
print(s1)    # hello world!!!hello world!!!
#endregion

# region    # ------------ 字符串的运算：比较运算 ------------
s1 = 'a whole new world'
s2 = 'hello world'
print(s1 == s2, s1 < s2)      # False True
print(s2 == 'hello world')    # True
print(s2 == 'Hello world')    # False
print(s2 != 'Hello world')    # True
s3 = '骆昊'
print(ord('骆'), ord('昊'))               # 39558 26122
s4 = '王大锤'
print(ord('王'), ord('大'), ord('锤'))    # 29579 22823 38180
print(s3 > s4, s3 <= s4)      # True False

s1 = 'hello world'
s2 = 'hello world'
s3 = s2
# 比较字符串的内容
print(s1 == s2, s2 == s3)    # True True
# 比较字符串的内存地址
print(s1 is s2, s2 is s3)    # False True
#endregion

# region    # ------------ 字符串的运算：成员运算 ------------
s1 = 'hello, world'
print('wo' in s1)    # True
s2 = 'goodbye'
print(s2 in s1)      # False
#endregion

# region    # ------------ 字符串的运算：获取字符串长度 ------------
s = 'hello, world'
print(len(s))                   # 12
print(len('goodbye, world'))    # 14
#endregion

# region    # ------------ 字符串的运算：索引 ------------
s = 'abc123456'
N = len(s)
# 获取第一个字符
print(s[0], s[-N])    # a a
# 获取最后一个字符
print(s[N-1], s[-1])  # 6 6
# 获取索引为2或-7的字符
print(s[2], s[-7])    # c c
# 获取索引为5和-4的字符
print(s[5], s[-4])    # 3 3
#endregion

# region    # ------------ 字符串的运算：切片 ------------
s = 'abc123456'
# i=2, j=5, k=1的正向切片操作
print(s[2:5])       # c12
# i=-7, j=-4, k=1的正向切片操作
print(s[-7:-4])     # c12
# i=2, j=9, k=1的正向切片操作
print(s[2:])        # c123456
# i=-7, j=9, k=1的正向切片操作
print(s[-7:])       # c123456
# i=2, j=9, k=2的正向切片操作
print(s[2::2])      # c246
# i=-7, j=9, k=2的正向切片操作
print(s[-7::2])     # c246
# i=0, j=9, k=2的正向切片操作
print(s[::2])       # ac246
# i=1, j=-1, k=2的正向切片操作
print(s[1:-1:2])    # b135
# i=7, j=1, k=-1的负向切片操作
print(s[7:1:-1])    # 54321c
# i=-2, j=-8, k=-1的负向切片操作
print(s[-2:-8:-1])  # 54321c
# i=7, j=-10, k=-1的负向切片操作
print(s[7::-1])     # 54321cba
# i=-1, j=1, k=-1的负向切片操作
print(s[:1:-1])     # 654321c
# i=0, j=9, k=1的正向切片
print(s[:])         # abc123456
# i=0, j=9, k=2的正向切片
print(s[::2])       # ac246
# i=-1, j=-10, k=-1的负向切片
print(s[::-1])      # 654321cba
# i=-1, j=-10, k=-2的负向切片
print(s[::-2])      # 642ca
#endregion

# region    # ------------ 字符串的运算：循环遍历每个字符 ------------
s1 = 'hello'
for index in range(len(s1)):
    print(s1[index])

s1 = 'hello'
for ch in s1:
    print(ch)
#endregion

# region    # ------------ 字符串方法：大小写相关操作 ------------
s1 = 'hello, world!'
# 使用capitalize方法获得字符串首字母大写后的字符串
print(s1.capitalize())   # Hello, world!
# 使用title方法获得字符串每个单词首字母大写后的字符串
print(s1.title())        # Hello, World!
# 使用upper方法获得字符串大写后的字符串
print(s1.upper())        # HELLO, WORLD!
s2 = 'GOODBYE'
# 使用lower方法获得字符串小写后的字符串
print(s2.lower())        # goodbye
#endregion

# region    # ------------ 字符串方法：查找操作 ------------
s = 'hello, world!'
# find方法从字符串中查找另一个字符串所在的位置
# 找到了返回字符串中另一个字符串首字符的索引
print(s.find('or'))        # 8
# 找不到返回-1
print(s.find('shit'))      # -1
# index方法与find方法类似
# 找到了返回字符串中另一个字符串首字符的索引
print(s.index('or'))       # 8
# 找不到引发异常
print(s.index('shit'))     # ValueError: substring not found
#endregion

# region    # ------------ 字符串方法：性质判断 ------------
s1 = 'hello, world!'
# startwith方法检查字符串是否以指定的字符串开头返回布尔值
print(s1.startswith('He'))    # False
print(s1.startswith('hel'))   # True
# endswith方法检查字符串是否以指定的字符串结尾返回布尔值
print(s1.endswith('!'))       # True

s2 = 'abc123456'
# isdigit方法检查字符串是否由数字构成返回布尔值
print(s2.isdigit())    # False
# isalpha方法检查字符串是否以字母构成返回布尔值
print(s2.isalpha())    # False
# isalnum方法检查字符串是否以数字和字母构成返回布尔值
print(s2.isalnum())    # True
#endregion

# region    # ------------ 字符串方法：格式化字符串 ------------
s = 'hello, world'
# center方法以宽度20将字符串居中并在两侧填充*
print(s.center(20, '*'))  # ****hello, world****
# rjust方法以宽度20将字符串右对齐并在左侧填充空格
print(s.rjust(20))        #         hello, world
# ljust方法以宽度20将字符串左对齐并在右侧填充~
print(s.ljust(20, '~'))   # hello, world~~~~~~~~
# 在字符串的左侧补零
print('33'.zfill(5))      # 00033
print('-33'.zfill(5))     # -0033

a = 321
b = 123
print('%d * %d = %d' % (a, b, a * b))

a = 321
b = 123
print('{0} * {1} = {2}'.format(a, b, a * b))

a = 321
b = 123
print(f'{a} * {b} = {a * b}')
#endregion

# region    # ------------ 字符串方法：修建操作 ------------
s = '   jackfrued@126.com  \t\r\n'
# strip方法获得字符串修剪左右两侧空格之后的字符串
print(s)    # jackfrued@126.com
print(s.strip())    # jackfrued@126.com
#endregion

# region    # ------------ 字符串方法：替换操作 ------------
s = 'hello, world'
print(s.replace('o', '@'))     # hell@, w@rld
print(s.replace('o', '@', 1))  # hell@, world
#endregion

# region    # ------------ 字符串方法：拆分/合并操作 ------------
s = 'I love you'
words = s.split()
print(words)            # ['I', 'love', 'you']
print('#'.join(words))  # I#love#you

s = 'I#love#you#so#much'
words = s.split('#')
print(words)  # ['I', 'love', 'you', 'so', 'much']
words = s.split('#', 3)
print(words)  # ['I', 'love', 'you', 'so#much']
#endregion

# region    # ------------ 字符串方法：编码/解码操作 ------------
a = '骆昊'
b = a.encode('utf-8')
c = a.encode('gbk')
print(b, c)  # b'\xe9\xaa\x86\xe6\x98\x8a' b'\xc2\xe6\xea\xbb'
print(b.decode('utf-8'))
print(c.decode('gbk'))
#endregion

# region    # ------------ 创建集合 ------------
# 创建集合的字面量语法(重复元素不会出现在集合中)
set1 = {1, 2, 3, 3, 3, 2}
print(set1)         # {1, 2, 3}
print(len(set1))    # 3
# 创建集合的构造器语法(后面会讲到什么是构造器)
set2 = set('hello')
print(set2)         # {'h', 'l', 'o', 'e'}
# 将列表转换成集合(可以去掉列表中的重复元素)
set3 = set([1, 2, 3, 3, 2, 1])
print(set3)         # {1, 2, 3}
# 创建集合的生成式语法(将列表生成式的[]换成{})
set4 = {num for num in range(1, 20) if num % 3 == 0 or num % 5 == 0}
print(set4)         # {3, 5, 6, 9, 10, 12, 15, 18}
# 集合元素的循环遍历
for elem in set4:
    print(elem)
#endregion

# region    # ------------ 集合的运算：成员运算 ------------
set1 = {11, 12, 13, 14, 15}
print(10 in set1)        # False
print(15 in set1)        # True
set2 = {'Python', 'Java', 'Go', 'Swift'}
print('Ruby' in set2)    # False
print('Java' in set2)    # True
#endregion

# region    # ------------ 集合的运算：交并差运算 ------------
set1 = {1, 2, 3, 4, 5, 6, 7}
set2 = {2, 4, 6, 8, 10}

# 交集
# 方法一: 使用 & 运算符
print(set1 & set2)                # {2, 4, 6}
# 方法二: 使用intersection方法
print(set1.intersection(set2))    # {2, 4, 6}

# 并集
# 方法一: 使用 | 运算符
print(set1 | set2)         # {1, 2, 3, 4, 5, 6, 7, 8, 10}
# 方法二: 使用union方法
print(set1.union(set2))    # {1, 2, 3, 4, 5, 6, 7, 8, 10}

# 差集
# 方法一: 使用 - 运算符
print(set1 - set2)              # {1, 3, 5, 7}
# 方法二: 使用difference方法
print(set1.difference(set2))    # {1, 3, 5, 7}

# 对称差
# 方法一: 使用 ^ 运算符
print(set1 ^ set2)                        # {1, 3, 5, 7, 8, 10}
# 方法二: 使用symmetric_difference方法
print(set1.symmetric_difference(set2))    # {1, 3, 5, 7, 8, 10}
# 方法三: 对称差相当于两个集合的并集减去交集
print((set1 | set2) - (set1 & set2))      # {1, 3, 5, 7, 8, 10}

# 复合赋值运算
set1 = {1, 3, 5, 7}
set2 = {2, 4, 6}
# 将set1和set2求并集再赋值给set1
# 也可以通过set1.update(set2)来实现
set1 |= set2
print(set1)    # {1, 2, 3, 4, 5, 6, 7}
set3 = {3, 6, 9}
# 将set1和set3求交集再赋值给set1
# 也可以通过set1.intersection_update(set3)来实现
set1 &= set3
print(set1)    # {3, 6}
#endregion

# region    # ------------ 集合的运算：比较运算 ------------
set1 = {1, 3, 5}
set2 = {1, 2, 3, 4, 5}
set3 = set2
# <运算符表示真子集，<=运算符表示子集
print(set1 < set2, set1 <= set2)    # True True
print(set2 < set3, set2 <= set3)    # False True
# 通过issubset方法也能进行子集判断
print(set1.issubset(set2))      # True

# 反过来可以用issuperset或>运算符进行超集判断
print(set2.issuperset(set1))    # True
print(set2 > set1)              # True
#endregion

# region    # ------------ 集合的方法 ------------
# 创建一个空集合
set1 = set()

# 通过add方法添加元素
set1.add(33)
set1.add(55)
set1.update({1, 10, 100, 1000})
print(set1)    # {33, 1, 100, 55, 1000, 10}

# 通过discard方法删除指定元素
set1.discard(100)
set1.discard(99)
print(set1)    # {1, 10, 33, 55, 1000}

# 通过remove方法删除指定元素，建议先做成员运算再删除
# 否则元素如果不在集合中就会引发KeyError异常
if 10 in set1:
    set1.remove(10)
print(set1)    # {33, 1, 55, 1000}

# pop方法可以从集合中随机删除一个元素并返回该元素
print(set1.pop())

# clear方法可以清空整个集合
set1.clear()
print(set1)    # set()

# 判断两个集合有没有相同的元素：有-F；没有-T
set1 = {'Java', 'Python', 'Go', 'Kotlin'}
set2 = {'Kotlin', 'Swift', 'Java', 'Objective-C', 'Dart'}
set3 = {'HTML', 'CSS', 'JavaScript'}
print(set1.isdisjoint(set2))    # False
print(set1.isdisjoint(set3))    # True
#endregion

# region    # ------------ 集合：不可变集合 ------------
set1 = frozenset({1, 3, 5, 7})
set2 = frozenset(range(1, 6))
print(set1 & set2)    # frozenset({1, 3, 5})
print(set1 | set2)    # frozenset({1, 2, 3, 4, 5, 7})
print(set1 - set2)    # frozenset({7})
print(set1 < set2)    # False
#endregion

# region    # ------------ 字典：创建和使用 ------------
### 方式一
xinhua = {
    '麓': '山脚下',
    '路': '道，往来通行的地方；方面，地区：南～货，外～货；种类：他俩是一～人',
    '蕗': '甘草的别名',
    '潞': '潞水，水名，即今山西省的浊漳河；潞江，水名，即云南省的怒江'
}
print(xinhua)
person = {
    'name': '王大锤', 'age': 55, 'weight': 60, 'office': '科华北路62号',
    'home': '中同仁路8号', 'tel': '13122334455', 'econtact': '13800998877'
}
print(person)

### 方式二：dict和zip
# dict函数(构造器)中的每一组参数就是字典中的一组键值对
person = dict(name='王大锤', age=55, weight=60, home='中同仁路8号')
print(person)    # {'name': '王大锤', 'age': 55, 'weight': 60, 'home': '中同仁路8号'}

# 可以通过Python内置函数zip压缩两个序列并创建字典
items1 = dict(zip('ABCDE', '12345'))
print(items1)    # {'A': '1', 'B': '2', 'C': '3', 'D': '4', 'E': '5'}
items2 = dict(zip('ABCDE', range(1, 10)))
print(items2)    # {'A': 1, 'B': 2, 'C': 3, 'D': 4, 'E': 5}

### 方式二：生成式
items3 = {x: x ** 3 for x in range(1, 6)}
print(items3)     # {1: 1, 2: 8, 3: 27, 4: 64, 5: 125}

### 获取键的值
person = {'name': '王大锤', 'age': 55, 'weight': 60, 'office': '科华北路62号'}
print(len(person))    # 4
for key in person:
    print(key)

for key in person:
    print(f'{key}: {person[key]}')
#endregion

# region    # ------------ 字典的运算 ------------
person = {'name': '王大锤', 'age': 55, 'weight': 60, 'office': '科华北路62号'}
# 检查name和tel两个键在不在person字典中
print('name' in person, 'tel' in person)    # True False
# 通过age修将person字典中对应的值修改为25
if 'age' in person:
    person['age'] = 25
# 通过索引操作向person字典中存入新的键值对
person['tel'] = '13122334455'
person['signature'] = 'xxxx，yyyyyy'
print('name' in person, 'tel' in person)    # True True
# 检查person字典中键值对的数量
print(len(person))    # 6
# 对字典的键进行循环并通索引运算获取键对应的值
for key in person:
    print(f'{key}: {person[key]}')
#endregion

# region    # ------------ 字典的方法 ------------
# 字典中的值又是一个字典(嵌套的字典)
students = {
    1001: {'name': '狄仁杰', 'sex': True, 'age': 22, 'place': '山西大同'},
    1002: {'name': '白元芳', 'sex': True, 'age': 23, 'place': '河北保定'},
    1003: {'name': '武则天', 'sex': False, 'age': 20, 'place': '四川广元'}
}

# 使用get方法通过键获取对应的值，如果取不到不会引发KeyError异常而是返回None或设定的默认值
print(students.get(1002))    # {'name': '白元芳', 'sex': True, 'age': 23, 'place': '河北保定'}
print(students.get(1005))    # None
print(students.get(1005, {'name': '无名氏'}))    # {'name': '无名氏'}

# 获取字典中所有的键
print(students.keys())      # dict_keys([1001, 1002, 1003])
# 获取字典中所有的值
print(students.values())    # dict_values([{...}, {...}, {...}])
# 获取字典中所有的键值对
print(students.items())     # dict_items([(1001, {...}), (1002, {....}), (1003, {...})])
# 对字典中所有的键值对进行循环遍历
for key, value in students.items():
    print(key, '--->', value)

# 使用pop方法通过键删除对应的键值对并返回该值
stu1 = students.pop(1002)
print(stu1)             # {'name': '白元芳', 'sex': True, 'age': 23, 'place': '河北保定'}
print(len(students))    # 2
# stu2 = students.pop(1005)    # KeyError: 1005
stu2 = students.pop(1005)
print(stu2)             # error
stu2 = students.pop(1005, {})
print(stu2)             # {}

# 使用popitem方法删除字典中最后一组键值对并返回对应的二元组
# 如果字典中没有元素，调用该方法将引发KeyError异常
key, value = students.popitem()
print(key, value)    # 1003 {'name': '武则天', 'sex': False, 'age': 20, 'place': '四川广元'}

# 如果这个键在字典中存在，setdefault返回原来与这个键对应的值
# 如果这个键在字典中不存在，向字典中添加键值对，返回第二个参数的值，默认为None
result = students.setdefault(1005, {'name': '方启鹤', 'sex': True})
print(result)        # {'name': '方启鹤', 'sex': True}
print(students)      # {1001: {...}, 1005: {...}}

# 使用update更新字典元素，相同的键会用新值覆盖掉旧值，不同的键会添加到字典中
others = {
    1005: {'name': '乔峰', 'sex': True, 'age': 32, 'place': '北京大兴'},
    1010: {'name': '王语嫣', 'sex': False, 'age': 19},
    1008: {'name': '钟灵', 'sex': False}
}
students.update(others)
print(students)      # {1001: {...}, 1005: {...}, 1010: {...}, 1008: {...}}

### 删除
person = {'name': '王大锤', 'age': 25, 'sex': True}
del person['age']
print(person)    # {'name': '王大锤', 'sex': True}
#endregion

# region    # ------------ 字典的应用：统计每个英文字母出现的次数 ------------
# sentence = input('请输入一段话: ')
sentence = 'aadddvv'
counter = {}
for ch in sentence:
    if 'A' <= ch <= 'Z' or 'a' <= ch <= 'z':
        counter[ch] = counter.get(ch, 0) + 1
for key, value in counter.items():
    print(f'字母{key}出现了{value}次.')
#endregion

# region    # ------------ 字典的应用：找出股价大于100元的股票并创建一个新的字典 ------------
stocks = {
    'AAPL': 191.88,
    'GOOG': 1186.96,
    'IBM': 149.24,
    'ORCL': 48.44,
    'ACN': 166.89,
    'FB': 208.09,
    'SYMC': 21.29
}
stocks2 = {key: value for key, value in stocks.items() if value > 100}
print(stocks2)
#endregion

# region    案例：输入M和N计算C(M,N)
def fac(num):
    """求阶乘"""
    result = 1
    for n in range(1, num + 1):
        result *= n
    # 返回num的阶乘（因变量）
    return result

# m = int(input('m = '))
# n = int(input('n = '))
m = int(7)
n = int(3)
# 当需要计算阶乘的时候不用再写重复的代码而是直接调用函数fac
# 调用函数的语法是在函数名后面跟上圆括号并传入参数
print(fac(m) // fac(n) // fac(m - n))
#endregion

# region    # ------------ 函数：可变参数 ------------
# 用星号表达式来表示args可以接收0个或任意多个参数
def add(*args):
    total = 0
    # 可变参数可以放在for循环中取出每个参数的值
    for val in args:
        if type(val) in (int, float):
            total += val
    return total

# 在调用add函数时可以传入0个或任意多个参数
print(add())
print(add(1))
print(add(1, 2))
print(add(1, 2, 3))
print(add(1, 3, 5, 7, 9))
#endregion

# region    # ------------ 函数：用模块管理函数 ------------
def foo():
    print('hello, world!')

def foo():
    print('goodbye, world!')

foo()  # 大家猜猜调用foo函数会输出什么
#endregion

# region    生成验证码
import random
import string

ALL_CHARS = string.digits + string.ascii_letters


def generate_code(code_len=4):
    """生成指定长度的验证码

    :param code_len: 验证码的长度(默认4个字符)
    :return: 由大小写英文字母和数字构成的随机验证码字符串
    """
    return ''.join(random.choices(ALL_CHARS, k=code_len))

for _ in range(10):
    print(generate_code())
#endregion

# region    返回给定文件的后缀名
### 方法一
def get_suffix(filename, ignore_dot=True):
    """获取文件名的后缀名

    :param filename: 文件名
    :param ignore_dot: 是否忽略后缀名前面的点
    :return: 文件的后缀名
    """
    # 从字符串中逆向查找.出现的位置
    pos = filename.rfind('.')
    # 通过切片操作从文件名中取出后缀名
    if pos <= 0:
        return ''
    return filename[pos + 1:] if ignore_dot else filename[pos:]

print(get_suffix('readme.txt'))       # txt
print(get_suffix('readme.txt.md'))    # md
print(get_suffix('.readme'))          #
print(get_suffix('readme.'))          #
print(get_suffix('readme'))           #

### 方法二
from os.path import splitext
def get_suffix(filename, ignore_dot=True):
    return splitext(filename)[1][1:]
    # return splitext(filename)[1]

print(get_suffix('readme.txt'))       # txt
print(get_suffix('readme.txt.md'))    # md
print(get_suffix('.readme'))          #
print(get_suffix('readme.'))          #
print(get_suffix('readme'))           #
#endregion

# region    判断给定的正整数是不是质数
def is_prime(num: int) -> bool:
    """判断一个正整数是不是质数

    :param num: 正整数
    :return: 如果是质数返回True，否则返回False
    """
    for i in range(2, int(num ** 0.5) + 1):
        if num % i == 0:
            return False
    return num != 1
#endregion

# region    计算两个正整数最大公约数和最小公倍数
### 方法一
def gcd_and_lcm(x: int, y: int) -> int:
    """求最大公约数和最小公倍数"""
    a, b = x, y
    while b % a != 0:
        a, b = b % a, a
    return a, x * y // a

### 方法二
def gcd(x: int, y: int) -> int:
    """求最大公约数"""
    while y % x != 0:
        x, y = y % x, x
    return x


def lcm(x: int, y: int) -> int:
    """求最小公倍数"""
    return x * y // gcd(x, y)
#endregion

# region    计算一组样本数据描述性统计信息
import math


def ptp(data):
    """求极差（全距）"""
    return max(data) - min(data)


def average(data):
    """求均值"""
    return sum(data) / len(data)


def variance(data):
    """求方差"""
    x_bar = average(data)
    temp = [(num - x_bar) ** 2 for num in data]
    return sum(temp) / (len(temp) - 1)


def standard_deviation(data):
    """求标准差"""
    return math.sqrt(variance(data))

def median(data):
    """找中位数"""
    temp, size = sorted(data), len(data)
    if size % 2 != 0:
        return temp[size // 2]
    else:
        return average(temp[size // 2 - 1:size // 2 + 1])
#endregion

# region    # ------------ 函数：关键字参数 ------------
def is_triangle(a, b, c):
    print(f'a = {a}, b = {b}, c = {c}')
    return a + b > c and b + c > a and a + c > b

# 调用函数传入参数不指定参数名按位置对号入座
print(is_triangle(1, 2, 3))
# 调用函数通过“参数名=参数值”的形式按顺序传入参数
print(is_triangle(a=1, b=2, c=3))
# 调用函数通过“参数名=参数值”的形式不按顺序传入参数
print(is_triangle(c=3, a=1, b=2))

def is_triangle(*, a, b, c):
    print(f'a = {a}, b = {b}, c = {c}')
    return a + b > c and b + c > a and a + c > b

# TypeError: is_triangle() takes 0 positional arguments but 3 were given
# print(is_triangle(3, 4, 5))
# 传参时必须使用“参数名=参数值”的方式，位置不重要
print(is_triangle(a=3, b=4, c=5))
print(is_triangle(c=5, b=4, a=3))

def calc(*args, **kwargs):
    result = 0
    for arg in args:
        if type(arg) in (int, float):
            result += arg
    for value in kwargs.values():
        if type(value) in (int, float):
            result += value
    return result

print(calc())                  # 0
print(calc(1, 2, 3))           # 6
print(calc(a=1, b=2, c=3))     # 6
print(calc(1, 2, c=3, d=4))    # 10
#endregion

# region    # ------------ 函数：装饰器 ------------
import random
import time

def download(filename):
    print(f'开始下载{filename}.')
    time.sleep(random.randint(2, 6))
    print(f'{filename}下载完成.')

def upload(filename):
    print(f'开始上传{filename}.')
    time.sleep(random.randint(4, 8))
    print(f'{filename}上传完成.')

download('MySQL从删库到跑路.avi')
upload('Python从入门到住院.pdf')

import time
# 定义装饰器函数，它的参数是被装饰的函数或类
def record_time(func):
    # 定义一个带装饰功能（记录被装饰函数的执行时间）的函数
    # 因为不知道被装饰的函数有怎样的参数所以使用*args和**kwargs接收所有参数
    # 在Python中函数可以嵌套的定义（函数中可以再定义函数）
    def wrapper(*args, **kwargs):
        # 在执行被装饰的函数之前记录开始时间
        start = time.time()
        # 执行被装饰的函数并获取返回值
        result = func(*args, **kwargs)
        # 在执行被装饰的函数之后记录结束时间
        end = time.time()
        # 计算和显示被装饰函数的执行时间
        print(f'{func.__name__}执行时间: {end - start:.3f}秒')
        # 返回被装饰函数的返回值（装饰器通常不会改变被装饰函数的执行结果）
        return result

    # 返回带装饰功能的wrapper函数
    return wrapper

download = record_time(download)
upload = record_time(upload)
download('MySQL从删库到跑路.avi')
upload('Python从入门到住院.pdf')
#endregion

# region    # ------------ 函数：语法糖 ------------
import random
import time

def record_time(func):

    def wrapper(*args, **kwargs):
        start = time.time()
        result = func(*args, **kwargs)
        end = time.time()
        print(f'{func.__name__}执行时间: {end - start:.3f}秒')
        return result

    return wrapper


@record_time
def download(filename):
    print(f'开始下载{filename}.')
    time.sleep(random.randint(2, 6))
    print(f'{filename}下载完成.')


@record_time
def upload(filename):
    print(f'开始上传{filename}.')
    time.sleep(random.randint(4, 8))
    print(f'{filename}上传完成.')


download('MySQL从删库到跑路.avi')
upload('Python从入门到住院.pdf')

#endregion

# region    # ------------ 函数：取消装饰器 ------------
import random
import time

from functools import wraps

def record_time(func):

    @wraps(func)
    def wrapper(*args, **kwargs):
        start = time.time()
        result = func(*args, **kwargs)
        end = time.time()
        print(f'{func.__name__}执行时间: {end - start:.3f}秒')
        return result

    return wrapper

@record_time
def download(filename):
    print(f'开始下载{filename}.')
    time.sleep(random.randint(2, 6))
    print(f'{filename}下载完成.')

@record_time
def upload(filename):
    print(f'开始上传{filename}.')
    time.sleep(random.randint(4, 8))
    print(f'{filename}上传完成.')

download('MySQL从删库到跑路.avi')
upload('Python从入门到住院.pdf')
# 取消装饰器
download.__wrapped__('MySQL必知必会.pdf')
upload = upload.__wrapped__
upload('Python从新手到大师.pdf')
#endregion

# region    # ------------ 函数：递归调用 ------------
### 阶乘
def fac(num):
    if num in (0, 1):
        return 1
    return num * fac(num - 1)

### 计算斐波那契数
def fib(n):
    if n in (1, 2):
        return 1
    return fib(n - 1) + fib(n - 2)

# 打印前20个斐波那契数
for i in range(1, 21):
    print(fib(i))

### 循环递推
def fib(n):
    a, b = 0, 1
    for _ in range(n):
        a, b = b, a + b
        # print(a, b)
    return a

for i in range(1, 21):
    print(fib(i))
#endregion

# region    # ------------ 类 ------------
### 定义类
class Student:

    def study(self, course_name):
        print(f'学生正在学习{course_name}.')

    def play(self):
        print(f'学生正在玩游戏.')

### 创建和使用对象
# 构造器：Student()
stu1 = Student()
stu2 = Student()
print(stu1)    # <__main__.Student object at 0x10ad5ac50>
print(stu2)    # <__main__.Student object at 0x10ad5acd0>
print(hex(id(stu1)), hex(id(stu2)))    # 0x10ad5ac50 0x10ad5acd0
# 所以stu3 = stu2这样的赋值语句并没有创建新的对象，只是用一个新的变量保存了已有对象的地址。

# 通过“类.方法”调用方法，第一个参数是接收消息的对象，第二个参数是学习的课程名称
Student.study(stu1, 'Python程序设计')    # 学生正在学习Python程序设计.
# 通过“对象.方法”调用方法，点前面的对象就是接收消息的对象，只需要传入第二个参数
stu1.study('Python程序设计')             # 学生正在学习Python程序设计.

Student.play(stu2)    # 学生正在玩游戏.
stu2.play()           # 学生正在玩游戏.
#endregion

# region    # ------------ 初始化方法 ------------
class Student:
    """学生"""

    def __init__(self, name, age):
        """初始化方法"""
        self.name = name
        self.age = age

    def study(self, course_name):
        """学习"""
        print(f'{self.name}正在学习{course_name}.')

    def play(self):
        """玩耍"""
        print(f'{self.name}正在玩游戏.')

# 由于初始化方法除了self之外还有两个参数
# 所以调用Student类的构造器创建对象时要传入这两个参数
stu1 = Student('骆昊', 40)
stu2 = Student('王大锤', 15)
stu2.name
stu2.age
stu1.study('Python程序设计')    # 骆昊正在学习Python程序设计.
stu2.play()                    # 王大锤正在玩游戏.
#endregion

# region    # ------------ 打印对象 ------------
class Student:
    """学生"""

    def __init__(self, name, age):
        """初始化方法"""
        self.name = name
        self.age = age

    def study(self, course_name):
        """学习"""
        print(f'{self.name}正在学习{course_name}.')

    def play(self):
        """玩耍"""
        print(f'{self.name}正在玩游戏.')

    def __repr__(self):
        return f'{self.name}: {self.age}'

stu1 = Student('骆昊', 40)
print(stu1)  # 骆昊: 40
students = [stu1, Student('李元芳', 36), Student('王大锤', 25)]
print(students)  # [骆昊: 40, 李元芳: 36, 王大锤: 25]
#endregion

# region    案例：定义一个类描述数字时钟。
import time

# 定义数字时钟类
class Clock(object):
    """数字时钟"""

    def __init__(self, hour=0, minute=0, second=0):
        """初始化方法
        :param hour: 时
        :param minute: 分
        :param second: 秒
        """
        self.hour = hour
        self.min = minute
        self.sec = second

    def run(self):
        """走字"""
        self.sec += 1
        if self.sec == 60:
            self.sec = 0
            self.min += 1
            if self.min == 60:
                self.min = 0
                self.hour += 1
                if self.hour == 24:
                    self.hour = 0

    def show(self):
        """显示时间"""
        return f'{self.hour:0>2d}:{self.min:0>2d}:{self.sec:0>2d}'


# 创建时钟对象
clock = Clock(23, 59, 58)
while True:
    # 给时钟对象发消息读取时间
    print(clock.show())
    # 休眠1秒钟
    time.sleep(1)
    # 给时钟对象发消息使其走字
    clock.run()
#endregion

# region    案例：定义一个类描述平面上的点，要求提供计算到另一个点距离的方法。
class Point(object):
    def __init__(self, x=0, y=0):
        """初始化方法
        :param x: 横坐标
        :param y: 纵坐标
        """
        self.x, self.y = x, y

    def __str__(self):
        return f'({self.x}, {self.y})'

    def distance_to(self, other):
        """计算与另一个点的距离
        :param other: 另一个点
        """
        dx = self.x - other.x
        dy = self.y - other.y
        return (dx * dx + dy * dy) ** 0.5


p1 = Point(3, 5)
p2 = Point(6, 9)
print(p1, p2)
print(p1.distance_to(p2))
print(p2.distance_to(p1))
#endregion

# region    # ------------ 可见性和属性装饰器 ------------
# __name表示一个私有属性，_name表示一个受保护属性
class Student:

    def __init__(self, name, age):
        self.__name = name
        self._age = age

    def study(self, course_name):
        print(f'{self.__name}正在学习{course_name}.')

stu = Student('王大锤', 20)
stu.study('Python程序设计')
print(stu.name)    # Error
print(stu.age)     # Error
print(stu.__name)  # Error
print(stu._age)    # Error

print(stu._Student__name)    # Right
print(stu._Student_age)      # Error

# 可以通过property装饰器为“私有”属性提供读取和修改的方法
class Student:

    def __init__(self, name, age):
        self.__name = name
        self.__age = age

    # 属性访问器(getter方法) - 获取__name属性
    @property
    def name(self):
        return self.__name

    # 属性修改器(setter方法) - 修改__name属性
    @name.setter
    def name(self, name):
        # 如果name参数不为空就赋值给对象的__name属性
        # 否则将__name属性赋值为'无名氏'，有两种写法
        # self.__name = name if name else '无名氏'
        self.__name = name or '无名氏'

    @property
    def age(self):
        return self.__age


stu = Student('王大锤', 20)
print(stu.name)  # 王大锤
print(stu.age)   #  20
stu.name = ''
print(stu.name)  # 无名氏
# stu.age = 30     # AttributeError: can't set attribute
#endregion

# region    # ------------ 动态属性 ------------
class Student:

    def __init__(self, name, age):
        self.name = name
        self.age = age


stu = Student('王大锤', 20)
# 为Student对象动态添加sex属性
stu.sex = '男'
print(stu.name)
print(stu.sex)

# 锁定属性
class Student:
    __slots__ = ('name', 'age')

    def __init__(self, name, age):
        self.name = name
        self.age = age


stu = Student('王大锤', 20)
# AttributeError: 'Student' object has no attribute 'sex'
stu.sex = '男'
#endregion

# region    # ------------ 静态方法和类方法 ------------
class Triangle(object):
    """三角形类"""

    def __init__(self, a, b, c):
        """初始化方法"""
        self.a = a
        self.b = b
        self.c = c

    @staticmethod
    def is_valid(a, b, c):
        """判断三条边长能否构成三角形(静态方法)"""
        return a + b > c and b + c > a and a + c > b

    # @classmethod
    # def is_valid(cls, a, b, c):
    #     """判断三条边长能否构成三角形(类方法)"""
    #     return a + b > c and b + c > a and a + c > b

    def perimeter(self):
        """计算周长"""
        return self.a + self.b + self.c

    def area(self):
        """计算面积"""
        p = self.perimeter() / 2
        return (p * (p - self.a) * (p - self.b) * (p - self.c)) ** 0.5
#endregion

# region    # ------------ 继承和多态 ------------
class Person:
    """人类"""

    def __init__(self, name, age):
        self.name = name
        self.age = age

    def eat(self):
        print(f'{self.name}正在吃饭.')

    def sleep(self):
        print(f'{self.name}正在睡觉.')


class Student(Person):
    """学生类"""

    def __init__(self, name, age):
        # super(Student, self).__init__(name, age)
        super().__init__(name, age)

    def study(self, course_name):
        print(f'{self.name}正在学习{course_name}.')


class Teacher(Person):
    """老师类"""

    def __init__(self, name, age, title):
        # super(Teacher, self).__init__(name, age)
        super().__init__(name, age)
        self.title = title

    def teach(self, course_name):
        print(f'{self.name}{self.title}正在讲授{course_name}.')


stu1 = Student('白元芳', 21)
stu2 = Student('狄仁杰', 22)
teacher = Teacher('武则天', 35, '副教授')
stu1.eat()
stu2.sleep()
teacher.teach('Python程序设计')
stu1.study('Python程序设计')
#endregion

# region    # ------------ 继承和多态 ------------
class Person:
    """人类"""

    def __init__(self, name, age):
        self.name = name
        self.age = age

    def eat(self):
        print(f'{self.name}正在吃饭.')

    def sleep(self):
        print(f'{self.name}正在睡觉.')


class Student(Person):
    """学生类"""

    def __init__(self, name, age):
        # super(Student, self).__init__(name, age)
        super().__init__(name, age)

    def study(self, course_name):
        print(f'{self.name}正在学习{course_name}.')

class Teacher(Person):
    """老师类"""

    def __init__(self, name, age, title):
        # super(Teacher, self).__init__(name, age)
        super().__init__(name, age)
        self.title = title

    def teach(self, course_name):
        print(f'{self.name}{self.title}正在讲授{course_name}.')

stu1 = Student('白元芳', 21)
stu2 = Student('狄仁杰', 22)
teacher = Teacher('武则天', 35, '副教授')
stu1.eat()
stu2.sleep()
teacher.teach('Python程序设计')
stu1.study('Python程序设计')
#endregion

# region    案例：扑克游戏
# ------------ 创建花色 ------------
from enum import Enum  # 创建枚举类型

class Suite(Enum):
    """花色(枚举)"""
    SPADE, HEART, CLUB, DIAMOND = range(4)

for suite in Suite:
    print(f'{suite}: {suite.value}')

# ------------ 定义牌类 ------------
class Card:
    """牌"""

    def __init__(self, suite, face):
        self.suite = suite
        self.face = face

    def __repr__(self):
        suites = '♠♥♣♦'
        faces = ['', 'A', '2', '3', '4', '5', '6', '7', '8', '9', '10', 'J', 'Q', 'K']
        # 根据牌的花色和点数取到对应的字符
        return f'{suites[self.suite.value]}{faces[self.face]}'

    def __lt__(self, other):
        # 花色相同比较点数的大小
        if self.suite == other.suite:
            return self.face < other.face
        # 花色不同比较花色对应的值
        return self.suite.value < other.suite.value

card1 = Card(Suite.SPADE, 5)
card2 = Card(Suite.HEART, 13)
print(card1, card2)    # ♠5 ♥K

# ------------ 定义牌类 ------------
import random

class Poker:
    """扑克"""

    def __init__(self):
        # 通过列表的生成式语法创建一个装52张牌的列表
        self.cards = [Card(suite, face) for suite in Suite
                      for face in range(1, 14)]
        # current属性表示发牌的位置
        self.current = 0

    def shuffle(self):
        """洗牌"""
        self.current = 0
        # 通过random模块的shuffle函数实现列表的随机乱序
        random.shuffle(self.cards)

    def deal(self):
        """发牌"""
        card = self.cards[self.current]
        self.current += 1
        return card

    @property
    def has_next(self):
        """还有没有牌可以发"""
        return self.current < len(self.cards)

poker = Poker()
poker.shuffle()
print(poker.cards)

# ------------ 定义玩家类 ------------
class Player:
    """玩家"""

    def __init__(self, name):
        self.name = name
        self.cards = []

    def get_one(self, card):
        """摸牌"""
        self.cards.append(card)

    def arrange(self):
        self.cards.sort()

# ------------ 创建四个玩家并将牌发到玩家的手上 ------------
poker = Poker()
poker.shuffle()
players = [Player('东邪'), Player('西毒'), Player('南帝'), Player('北丐')]

for _ in range(13):
    for player in players:
        player.get_one(poker.deal())

for player in players:
    player.arrange()
    print(f'{player.name}: ', end='')
    print(player.cards)

#endregion

# region    案例：工资结算系统
from abc import ABCMeta, abstractmethod  # 用ABCMeta元类来定义抽象类

class Employee(metaclass=ABCMeta):
    """员工"""

    def __init__(self, name):
        self.name = name

    @abstractmethod
    def get_salary(self):
        """结算月薪"""
        pass

class Manager(Employee):
    """部门经理"""

    def get_salary(self):
        return 15000.0

class Programmer(Employee):
    """程序员"""

    def __init__(self, name, working_hour=0):
        super().__init__(name)
        self.working_hour = working_hour

    def get_salary(self):
        return 200 * self.working_hour

class Salesman(Employee):
    """销售员"""

    def __init__(self, name, sales=0):
        super().__init__(name)
        self.sales = sales

    def get_salary(self):
        return 1800 + self.sales * 0.05

emps = [
    Manager('刘备'), Programmer('诸葛亮'), Manager('曹操'),
    Programmer('荀彧'), Salesman('吕布'), Programmer('张辽'),
]
for emp in emps:
    if isinstance(emp, Programmer):
        emp.working_hour = int(input(f'请输入{emp.name}本月工作时间: '))
    elif isinstance(emp, Salesman):
        emp.sales = float(input(f'请输入{emp.name}本月销售额: '))
    print(f'{emp.name}本月工资为: ￥{emp.get_salary():.2f}元')

#endregion

# region    # ------------ base64 - Base64编解码模块 ------------
import base64

content = 'Man is distinguished, not only by his reason, but by this singular passion from other animals, which is a lust of the mind, that by a perseverance of delight in the continued and indefatigable generation of knowledge, exceeds the short vehemence of any carnal pleasure.'
base64.b64encode(content.encode())
# b'TWFuIGlzIGRpc3Rpbmd1aXNoZWQsIG5vdCBvbmx5IGJ5IGhpcyByZWFzb24sIGJ1dCBieSB0aGlzIHNpbmd1bGFyIHBhc3Npb24gZnJvbSBvdGhlciBhbmltYWxzLCB3aGljaCBpcyBhIGx1c3Qgb2YgdGhlIG1pbmQsIHRoYXQgYnkgYSBwZXJzZXZlcmFuY2Ugb2YgZGVsaWdodCBpbiB0aGUgY29udGludWVkIGFuZCBpbmRlZmF0aWdhYmxlIGdlbmVyYXRpb24gb2Yga25vd2xlZGdlLCBleGNlZWRzIHRoZSBzaG9ydCB2ZWhlbWVuY2Ugb2YgYW55IGNhcm5hbCBwbGVhc3VyZS4='
content = b'TWFuIGlzIGRpc3Rpbmd1aXNoZWQsIG5vdCBvbmx5IGJ5IGhpcyByZWFzb24sIGJ1dCBieSB0aGlzIHNpbmd1bGFyIHBhc3Npb24gZnJvbSBvdGhlciBhbmltYWxzLCB3aGljaCBpcyBhIGx1c3Qgb2YgdGhlIG1pbmQsIHRoYXQgYnkgYSBwZXJzZXZlcmFuY2Ugb2YgZGVsaWdodCBpbiB0aGUgY29udGludWVkIGFuZCBpbmRlZmF0aWdhYmxlIGdlbmVyYXRpb24gb2Yga25vd2xlZGdlLCBleGNlZWRzIHRoZSBzaG9ydCB2ZWhlbWVuY2Ugb2YgYW55IGNhcm5hbCBwbGVhc3VyZS4='
base64.b64decode(content).decode()
# 'Man is distinguished, not only by his reason, but by this singular passion from other animals, which is a lust of the mind, that by a perseverance of delight in the continued and indefatigable generation of knowledge, exceeds the short vehemence of any carnal pleasure.'
#endregion

# region    # ------------ collections - 容器数据类型模块 ------------
### 使用namedtuple创建扑克牌类
from collections import namedtuple

Card = namedtuple('Card', ('suite', 'face'))
card1 = Card('红桃', 5)
card2 = Card('草花', 9)

print(f'{card1.suite}{card1.face}')
print(f'{card2.suite}{card2.face}')

### 使用Counter类统计列表中出现次数最多的三个元素
from collections import Counter

words = [
    'look', 'into', 'my', 'eyes', 'look', 'into', 'my', 'eyes',
    'the', 'eyes', 'the', 'eyes', 'the', 'eyes', 'not', 'around',
    'the', 'eyes', "don't", 'look', 'around', 'the', 'eyes',
    'look', 'into', 'my', 'eyes', "you're", 'under'
]
counter = Counter(words)
# 打印words列表中出现频率最高的3个元素及其出现次数
for elem, count in counter.most_common(3):
    print(elem, count)
#endregion

# region    案例：hashlib - 哈希函数模块
import hashlib

# 计算字符串"123456"的MD5摘要
print(hashlib.md5('123456'.encode()).hexdigest())

# 计算文件"Python-3.7.1.tar.xz"的MD5摘要
hasher = hashlib.md5()
with open('Python-3.7.1.tar.xz', 'rb') as file:
    data = file.read(512)
    while data:
        hasher.update(data)
        data = file.read(512)
print(hasher.hexdigest())
#endregion

# region    案例：heapq - 堆排序模块
import heapq

list1 = [34, 25, 12, 99, 87, 63, 58, 78, 88, 92]
# 找出列表中最大的三个元素
print(heapq.nlargest(3, list1))
# 找出列表中最小的三个元素
print(heapq.nsmallest(3, list1))

list2 = [
    {'name': 'IBM', 'shares': 100, 'price': 91.1},
    {'name': 'AAPL', 'shares': 50, 'price': 543.22},
    {'name': 'FB', 'shares': 200, 'price': 21.09},
    {'name': 'HPQ', 'shares': 35, 'price': 31.75},
    {'name': 'YHOO', 'shares': 45, 'price': 16.35},
    {'name': 'ACME', 'shares': 75, 'price': 115.65}
]
# 找出价格最高的三只股票
print(heapq.nlargest(3, list2, key=lambda x: x['price']))
# 找出持有数量最高的三只股票
print(heapq.nlargest(3, list2, key=lambda x: x['shares']))
#endregion

# region    案例：itertools - 迭代工具模块
import itertools

# 产生ABCD的全排列
for value in itertools.permutations('ABCD'):
    print(value)

# 产生ABCDE的五选三组合
for value in itertools.combinations('ABCDE', 3):
    print(value)

# 产生ABCD和123的笛卡尔积
for value in itertools.product('ABCD', '123'):
    print(value)

# 产生ABC的无限循环序列
it = itertools.cycle(('A', 'B', 'C'))
print(next(it))
print(next(it))
print(next(it))
print(next(it))
#endregion

# region    案例：uuid - UUID生成模块
import uuid

uuid.uuid1().hex
uuid.uuid1().hex
uuid.uuid1().hex
#endregion

# region    # ------------ 文件读取 ------------
### 整个文件读取
file = open('致橡树.txt', 'r', encoding='utf-8')
print(file.read())
file.close()

### 按行读取
file = open('致橡树.txt', 'r', encoding='utf-8')
for line in file:
    print(line, end='')
file.close()

file = open('致橡树.txt', 'r', encoding='utf-8')
lines = file.readlines()
for line in lines:
    print(line, end='')
file.close()

### 追加写入
file = open('致橡树.txt', 'a', encoding='utf-8')
file.write('\n标题：《致橡树》')
file.write('\n作者：舒婷')
file.write('\n时间：1977年3月')
file.close()
#endregion

# region    # ------------ 异常处理机制 ------------
file = None
try:
    file = open('致橡树.txt', 'r', encoding='utf-8')
    print(file.read())
except FileNotFoundError:
    print('无法打开指定的文件!')
except LookupError:
    print('指定了未知的编码!')
except UnicodeDecodeError:
    print('读取文件时解码错误!')
finally:
    if file:
        file.close()

### 案例：调用求阶乘的函数fac，通过try...except...结构捕获输入错误的异常并打印异常对象（显示异常信息），如果输入正确就计算阶乘并结束程序。
class InputError(ValueError):
    """自定义异常类型"""
    pass

def fac(num):
    """求阶乘"""
    if num < 0:
        raise InputError('只能计算非负整数的阶乘')
    if num in (0, 1):
        return 1
    return num * fac(num - 1)

flag = True
while flag:
    num = int(input('n = '))
    try:
        print(f'{num}! = {fac(num)}')
        flag = False
    except InputError as err:
        print(err)

#endregion

# region    # ------------ 上下文语法：读文件 ------------
try:
    with open('致橡树.txt', 'r', encoding='utf-8') as file:
        print(file.read())
except FileNotFoundError:
    print('无法打开指定的文件!')
except LookupError:
    print('指定了未知的编码!')
except UnicodeDecodeError:
    print('读取文件时解码错误!')
#endregion

# region    # ------------ 读写二进制文件 ------------
try:
    with open('guido.jpg', 'rb') as file1:
        data = file1.read()
    with open('吉多.jpg', 'wb') as file2:
        file2.write(data)
except FileNotFoundError:
    print('指定的文件无法打开.')
except IOError:
    print('读写文件时出现错误.')
print('程序执行结束.')
#endregion

# region    # ------------ 读写文件：指定size以节省内存开销 ------------
try:
    with open('guido.jpg', 'rb') as file1, open('吉多.jpg', 'wb') as file2:
        data = file1.read(512)
        while data:
            file2.write(data)
            data = file1.read()
except FileNotFoundError:
    print('指定的文件无法打开.')
except IOError:
    print('读写文件时出现错误.')
print('程序执行结束.')
#endregion

# region    # ------------ 生成word文档 ------------
from docx import Document
from docx.shared import Cm, Pt

from docx.document import Document as Doc

# 创建代表Word文档的Doc对象
document = Document()  # type: Doc
# 添加大标题
document.add_heading('快快乐乐学Python', 0)
# 添加段落
p = document.add_paragraph('Python是一门非常流行的编程语言，它')
run = p.add_run('简单')
run.bold = True
run.font.size = Pt(18)
p.add_run('而且')
run = p.add_run('优雅')
run.font.size = Pt(18)
run.underline = True
p.add_run('。')

# 添加一级标题
document.add_heading('Heading, level 1', level=1)
# 添加带样式的段落
document.add_paragraph('Intense quote', style='Intense Quote')
# 添加无序列表
document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'second item in ordered list', style='List Bullet'
)
# 添加有序列表
document.add_paragraph(
    'first item in ordered list', style='List Number'
)
document.add_paragraph(
    'second item in ordered list', style='List Number'
)

# 添加图片（注意路径和图片必须要存在）
document.add_picture('知乎专栏\\专栏资源\\第026课\\guido.jpg', width=Cm(5.2))

# 添加分节符
document.add_section()

records = (
    ('骆昊', '男', '1995-5-5'),
    ('孙美丽', '女', '1992-2-2')
)
# 添加表格
table = document.add_table(rows=1, cols=3)
table.style = 'Dark List'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = '姓名'
hdr_cells[1].text = '性别'
hdr_cells[2].text = '出生日期'
# 为表格添加行
for name, sex, birthday in records:
    row_cells = table.add_row().cells
    row_cells[0].text = name
    row_cells[1].text = sex
    row_cells[2].text = birthday

# 添加分页符
document.add_page_break()

# 保存文档
document.save('demo.docx')
#endregion

# region    # ------------ 读取Word ------------
from docx import Document
from docx.document import Document as Doc

doc = Document('知乎专栏\\专栏资源\\第026课\\离职证明.docx')  # type: Doc
for no, p in enumerate(doc.paragraphs):
    print(no, p.text)
#endregion

# region    # ------------ 基于模板批量生成word ------------
from docx import Document
from docx.document import Document as Doc

# 将真实信息用字典的方式保存在列表中
employees = [
    {
        'name': '骆昊',
        'id': '100200198011280001',
        'sdate': '2008年3月1日',
        'edate': '2012年2月29日',
        'department': '产品研发',
        'position': '架构师',
        'company': '成都华为技术有限公司'
    },
    {
        'name': '王大锤',
        'id': '510210199012125566',
        'sdate': '2019年1月1日',
        'edate': '2021年4月30日',
        'department': '产品研发',
        'position': 'Python开发工程师',
        'company': '成都谷道科技有限公司'
    },
    {
        'name': '李元芳',
        'id': '2102101995103221599',
        'sdate': '2020年5月10日',
        'edate': '2021年3月5日',
        'department': '产品研发',
        'position': 'Java开发工程师',
        'company': '同城企业管理集团有限公司'
    },
]
# 对列表进行循环遍历，批量生成Word文档
for emp_dict in employees:
    # 读取离职证明模板文件
    doc = Document('知乎专栏\\专栏资源\\第026课\\离职证明模板.docx')  # type: Doc
    # 循环遍历所有段落寻找占位符
    for p in doc.paragraphs:
        if '{' not in p.text:
            continue
        # 不能直接修改段落内容，否则会丢失样式
        # 所以需要对段落中的元素进行遍历并进行查找替换
        for run in p.runs:
            if '{' not in run.text:
                continue
            # 将占位符换成实际内容
            start, end = run.text.find('{'), run.text.find('}')
            key, place_holder = run.text[start + 1:end], run.text[start:end + 1]
            run.text = run.text.replace(place_holder, emp_dict[key])
    # 每个人对应保存一个Word文档
    doc.save(f'{emp_dict["name"]}离职证明.docx')
#endregion

# region    # ------------ 生成PowerPoint ------------
from pptx import Presentation

# 创建幻灯片对象
pres = Presentation()

# 选择母版添加一页
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)
# 获取标题栏和副标题栏
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 编辑标题和副标题
title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

# 选择母版添加一页
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)
# 获取页面上所有形状
shapes = slide.shapes
# 获取标题和主体
title_shape = shapes.title
body_shape = shapes.placeholders[1]
# 编辑标题
title_shape.text = 'Introduction'
# 编辑主体内容
tf = body_shape.text_frame
tf.text = 'History of Python'
# 添加一个一级段落
p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 1
# 添加一个二级段落
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python.'
p.level = 2

# 保存幻灯片
pres.save('test.pptx')
#endregion

# region    # ------------ 从PDF中提取文本 ------------
import PyPDF2

reader = PyPDF2.PdfReader('test.pdf')
for page in reader.pages:
    print(page.extract_text())
#endregion

# region    # ------------ 旋转和叠加页面 ------------
reader = PyPDF2.PdfReader('XGBoost.pdf')
writer = PyPDF2.PdfWriter()

for no, page in enumerate(reader.pages):
    if no % 2 == 0:
        new_page = page.rotate(-90)
    else:
        new_page = page.rotate(90)
    writer.add_page(new_page)

with open('temp.pdf', 'wb') as file_obj:
    writer.write(file_obj)
#endregion

# region    # ------------ 加密PDF文件 ------------
import PyPDF2

reader = PyPDF2.PdfReader('XGBoost.pdf')
writer = PyPDF2.PdfWriter()

for page in reader.pages:
    writer.add_page(page)

writer.encrypt('foobared')

with open('temp.pdf', 'wb') as file_obj:
    writer.write(file_obj)
#endregion

# region    # ------------ 批量添加水印 ------------
import PyPDF2

reader = PyPDF2.PdfReader('XGBoost.pdf')
writer = PyPDF2.PdfWriter()

for page in reader.pages:
    writer.add_page(page)

writer.encrypt('foobared')

with open('temp.pdf', 'wb') as file_obj:
    writer.write(file_obj)
#endregion

# region    # ------------ 创建PDF文件 ------------
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

pdf_canvas = canvas.Canvas('resources/demo.pdf', pagesize=A4)
width, height = A4

# 绘图
image = canvas.ImageReader('resources/guido.jpg')
pdf_canvas.drawImage(image, 20, height - 395, 250, 375)

# 显示当前页
pdf_canvas.showPage()

# 注册字体文件
pdfmetrics.registerFont(TTFont('Font1', 'resources/fonts/Vera.ttf'))
pdfmetrics.registerFont(TTFont('Font2', 'resources/fonts/青呱石头体.ttf'))

# 写字
pdf_canvas.setFont('Font2', 40)
pdf_canvas.setFillColorRGB(0.9, 0.5, 0.3, 1)
pdf_canvas.drawString(width // 2 - 120, height // 2, '你好，世界！')
pdf_canvas.setFont('Font1', 40)
pdf_canvas.setFillColorRGB(0, 1, 0, 0.5)
pdf_canvas.rotate(18)
pdf_canvas.drawString(250, 250, 'hello, world!')

# 保存
pdf_canvas.save()
#endregion

# region    # ------------ 用Pillow处理图像：读取和显示图像 ------------
from PIL import Image

# 读取图像获得Image对象
image = Image.open('知乎专栏\\专栏资源\\第026课\\guido.jpg')
# 通过Image对象的format属性获得图像的格式
print(image.format) # JPEG
# 通过Image对象的size属性获得图像的尺寸
print(image.size)   # (500, 750)
# 通过Image对象的mode属性获取图像的模式
print(image.mode)   # RGB
# 通过Image对象的show方法显示图像
image.show()
#endregion

# region    # ------------ 用Pillow处理图像：剪裁图像 ------------
# 通过Image对象的crop方法指定剪裁区域剪裁图像
image.crop((80, 20, 310, 360)).show()
#endregion

# region    # ------------ 用Pillow处理图像：生成缩略图 ------------
# 通过Image对象的thumbnail方法生成指定尺寸的缩略图
image.thumbnail((128, 128))
image.show()
#endregion

# region    # ------------ 用Pillow处理图像：缩放和黏贴图像 ------------
# 读取骆昊的照片获得Image对象
luohao_image = Image.open('luohao.png')
# 读取吉多的照片获得Image对象
guido_image = Image.open('guido.jpg')
# 从吉多的照片上剪裁出吉多的头
guido_head = guido_image.crop((80, 20, 310, 360))
width, height = guido_head.size
# 使用Image对象的resize方法修改图像的尺寸
# 使用Image对象的paste方法将吉多的头粘贴到骆昊的照片上
luohao_image.paste(guido_head.resize((int(width / 1.5), int(height / 1.5))), (172, 40))
luohao_image.show()
#endregion

# region    # ------------ 用Pillow处理图像：旋转和翻转 ------------
image = Image.open('guido.jpg')
# 使用Image对象的rotate方法实现图像的旋转
image.rotate(45).show()
# 使用Image对象的transpose方法实现图像翻转
# Image.FLIP_LEFT_RIGHT - 水平翻转
# Image.FLIP_TOP_BOTTOM - 垂直翻转
image.transpose(Image.FLIP_TOP_BOTTOM).show()
#endregion

# region    # ------------ 用Pillow处理图像：操作像素 ------------
for x in range(80, 310):
    for y in range(20, 360):
        # 通过Image对象的putpixel方法修改图像指定像素点
        image.putpixel((x, y), (128, 128, 128))
image.show()
#endregion

# region    # ------------ 用Pillow处理图像：滤镜效果 ------------
from PIL import ImageFilter

# 使用Image对象的filter方法对图像进行滤镜处理
# ImageFilter模块包含了诸多预设的滤镜也可以自定义滤镜
image.filter(ImageFilter.CONTOUR).show()
#endregion

# region    # ------------ 使用Pillow绘图 ------------
import random
from PIL import Image, ImageDraw, ImageFont

def random_color():
    """生成随机颜色"""
    red = random.randint(0, 255)
    green = random.randint(0, 255)
    blue = random.randint(0, 255)
    return red, green, blue


width, height = 800, 600
# 创建一个800*600的图像，背景色为白色
image = Image.new(mode='RGB', size=(width, height), color=(255, 255, 255))
# 创建一个ImageDraw对象
drawer = ImageDraw.Draw(image)
# 通过指定字体和大小获得ImageFont对象
font = ImageFont.truetype('Kongxin.ttf', 32)
# 通过ImageDraw对象的text方法绘制文字
drawer.text((300, 50), 'Hello, world!', fill=(255, 0, 0), font=font)
# 通过ImageDraw对象的line方法绘制两条对角直线
drawer.line((0, 0, width, height), fill=(0, 0, 255), width=2)
drawer.line((width, 0, 0, height), fill=(0, 0, 255), width=2)
xy = width // 2 - 60, height // 2 - 60, width // 2 + 60, height // 2 + 60
# 通过ImageDraw对象的rectangle方法绘制矩形
drawer.rectangle(xy, outline=(255, 0, 0), width=2)
# 通过ImageDraw对象的ellipse方法绘制椭圆
for i in range(4):
    left, top, right, bottom = 150 + i * 120, 220, 310 + i * 120, 380
    drawer.ellipse((left, top, right, bottom), outline=random_color(), width=8)
# 显示图像
image.show()
# 保存图像
image.save('result.png')
#endregion

# region    # ------------ 用Python发送邮件和短信 ------------
import smtplib
from email.header impOort Header
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

# 创建邮件主体对象
email = MIMEMultipart()
# 设置发件人、收件人和主题
email['From'] = 'xxxxxxxxx@126.com'
email['To'] = 'yyyyyy@qq.com;zzzzzz@1000phone.com'
email['Subject'] = Header('上半年工作情况汇报', 'utf-8')
# 添加邮件正文内容
content = """据德国媒体报道，当地时间9日，德国火车司机工会成员进行了投票，
定于当地时间10日起进行全国性罢工，货运交通方面的罢工已于当地时间10日19时开始。
此后，从11日凌晨2时到13日凌晨2时，德国全国范围内的客运和铁路基础设施将进行48小时的罢工。"""
email.attach(MIMEText(content, 'plain', 'utf-8'))

# 创建SMTP_SSL对象（连接邮件服务器）
smtp_obj = smtplib.SMTP_SSL('smtp.126.com', 465)
# 通过用户名和授权码进行登录
smtp_obj.login('xxxxxxxxx@126.com', '邮件服务器的授权码')
# 发送邮件（发件人、收件人、邮件内容（字符串））
smtp_obj.sendmail(
    'xxxxxxxxx@126.com',
    ['yyyyyy@qq.com', 'zzzzzz@1000phone.com'],
    email.as_string()
)
#endregion

# region    # ------------ 发送短信 ------------
import random
import requests

def send_message_by_luosimao(tel, message):
    """发送短信（调用螺丝帽短信网关）"""
    resp = requests.post(
        url='http://sms-api.luosimao.com/v1/send.json',
        auth=('api', 'key-注册成功后平台分配的KEY'),
        data={
            'mobile': tel,
            'message': message
        },
        timeout=10,
        verify=False
    )
    return resp.json()

def gen_mobile_code(length=6):
    """生成指定长度的手机验证码"""
    return ''.join(random.choices('0123456789', k=length))

def main():
    code = gen_mobile_code()
    message = f'您的短信验证码是{code}，打死也不能告诉别人哟！【Python小课】'
    print(send_message_by_luosimao('13500112233', message))

if __name__ == '__main__':
    main()
#endregion

# region    案例：验证输入用户名和QQ号是否有效并给出对应的提示信息。
"""
要求：用户名必须由字母、数字或下划线构成且长度在6~20个字符之间，QQ号是5~12的数字且首位不能为0
"""
import re

username = input('请输入用户名: ')
qq = input('请输入QQ号: ')
# match函数的第一个参数是正则表达式字符串或正则表达式对象
# match函数的第二个参数是要跟正则表达式做匹配的字符串对象
m1 = re.match(r'^[0-9a-zA-Z_]{6,20}$', username)
if not m1:
    print('请输入有效的用户名.')
# fullmatch函数要求字符串和正则表达式完全匹配
# 所以正则表达式没有写起始符和结束符
m2 = re.fullmatch(r'[1-9]\d{4,11}', qq)
if not m2:
    print('请输入有效的QQ号.')
if m1 and m2:
    print('你输入的信息是有效的!')
#endregion

# region    案例：从一段文字中提取出国内手机号码。
import re

# 创建正则表达式对象，使用了前瞻和回顾来保证手机号前后不应该再出现数字
pattern = re.compile(r'(?<=\D)1[34578]\d{9}(?=\D)')
sentence = '''重要的事情说8130123456789遍，我的手机号是13512346789这个靓号，
不是15600998765，也是110或119，王大锤的手机号才是15600998765。'''
# 方法一：查找所有匹配并保存到一个列表中
tels_list = re.findall(pattern, sentence)
for tel in tels_list:
    print(tel)
print('--------华丽的分隔线--------')

# 方法二：通过迭代器取出匹配对象并获得匹配的内容
for temp in pattern.finditer(sentence):
    print(temp.group())
print('--------华丽的分隔线--------')

# 方法三：通过search函数指定搜索位置找出所有匹配
m = pattern.search(sentence)
while m:
    print(m.group())
    m = pattern.search(sentence, m.end())
#endregion

# region    案例：替换字符串中的不良内容
import re

sentence = 'Oh, shit! 你是傻逼吗? Fuck you.'
purified = re.sub('fuck|shit|[傻煞沙][比笔逼叉缺吊碉雕]',
                  '*', sentence, flags=re.IGNORECASE)
print(purified)  # Oh, *! 你是*吗? * you.
#endregion

# region    案例：拆分长字符串
import re

poem = '窗前明月光，疑是地上霜。举头望明月，低头思故乡。'
sentences_list = re.split(r'[，。]', poem)
sentences_list = [sentence for sentence in sentences_list if sentence]
for sentence in sentences_list:
    print(sentence)
#endregion

# region    # ------------ 爬虫：获取搜狐首页代码 ------------
import requests

resp = requests.get('https://www.sohu.com/')
if resp.status_code == 200:
    print(resp.text)
#endregion

# region    # ------------ 爬虫：从页面的 HTML 代码中提取新闻的标题和链接 ------------
import re
import requests

pattern = re.compile(r'<a.*?href="(.*?)".*?title="(.*?)".*?>')
resp = requests.get('https://www.sohu.com/')
if resp.status_code == 200:
    all_matches = pattern.findall(resp.text)
    for href, title in all_matches:
        print(href)
        print(title)
#endregion

# region    # ------------ 爬虫：获取logo ------------
import requests

resp = requests.get('https://www.baidu.com/img/PCtm_d9c8750bed0b3c7d089fa7d55720d6cf.png')
with open('baidu.png', 'wb') as file:
    file.write(resp.content)
#endregion

# region    案例：爬取豆瓣电影
import random
import re
import time
import requests

for page in range(1, 11):
    resp = requests.get(
        url=f'https://movie.douban.com/top250?start={(page - 1) * 25}',
        # 如果不设置HTTP请求头中的User-Agent，豆瓣会检测出不是浏览器而阻止我们的请求。
        # 通过get函数的headers参数设置User-Agent的值，具体的值可以在浏览器的开发者工具查看到。
        # 用爬虫访问大部分网站时，将爬虫伪装成来自浏览器的请求都是非常重要的一步。
        headers={'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_14_6) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/92.0.4515.159 Safari/537.36'}
    )
    # 通过正则表达式获取class属性为title且标签体不以&开头的span标签并用捕获组提取标签内容
    pattern1 = re.compile(r'<span class="title">([^&]*?)</span>')
    titles = pattern1.findall(resp.text)
    # 通过正则表达式获取class属性为rating_num的span标签并用捕获组提取标签内容
    pattern2 = re.compile(r'<span class="rating_num".*?>(.*?)</span>')
    ranks = pattern2.findall(resp.text)
    # 使用zip压缩两个列表，循环遍历所有的电影标题和评分
    for title, rank in zip(titles, ranks):
        print(title, rank)
    # 随机休眠1-5秒，避免爬取页面过于频繁
    time.sleep(random.random() * 4 + 1)
#endregion











